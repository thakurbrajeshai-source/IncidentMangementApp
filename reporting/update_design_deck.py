"""
Adds a "Reporter-driven workflow trigger" slide to the Design Deck so the
deck stays in sync with the PRD after the reporter-side trigger change.

Run from repo root:
    python src/reporting/update_design_deck.py
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

DECK = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "..", "..", "Incident_App_Design_Deck.pptx",
)
DECK = os.path.normpath(DECK)

NAVY = RGBColor(0x1E, 0x3A, 0x8A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
MUTED = RGBColor(0x6B, 0x72, 0x80)
TEXT = RGBColor(0x11, 0x18, 0x27)
ACCENT_BG = RGBColor(0xEC, 0xFE, 0xFF)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
AMBER = RGBColor(0xD9, 0x77, 0x06)


def _add_text(slide, left, top, width, height, text, *, size=14, bold=False,
              color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_rect(slide, left, top, width, height, *, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def add_slide(prs):
    # Pick the first available layout (some decks only expose a few).
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # Title
    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.5),
              "Reporter-driven workflow trigger", size=26, bold=True, color=NAVY)
    _add_text(slide, Inches(0.5), Inches(0.85), Inches(12), Inches(0.4),
              "Reporters run the default check for their own ticket — and attach more inside the conversation",
              size=14, color=MUTED)

    # Left column: flow
    flow_x = Inches(0.5)
    flow_y = Inches(1.5)
    flow_w = Inches(7.2)
    step_h = Inches(0.7)

    _add_text(slide, flow_x, flow_y, flow_w, Inches(0.4), "Flow", size=16, bold=True, color=NAVY)
    steps = [
        ("1.", "Resolver/Admin marks a workflow as the default check for one or more categories"),
        ("2.", "Reporter creates an incident in that category — the workflow card auto-appears"),
        ("3.", "Reporter fills the required input (e.g. Lead ID) and hits Run"),
        ("4.", "Rendered result table renders inside the same card"),
        ("5.", "Reporter closes the ticket if satisfied — or replies that it's still broken"),
        ("6.", "Inside the conversation, the reporter or admin can Attach another check"),
    ]
    for i, (n, text) in enumerate(steps):
        y = flow_y + Inches(0.5) + Inches(0.7) * i
        _add_text(slide, flow_x, y, Inches(0.4), step_h, n, size=13, bold=True, color=TEAL)
        _add_text(slide, flow_x + Inches(0.4), y, flow_w - Inches(0.4), step_h, text,
                  size=13, color=TEXT)

    # Right column: wireframe mockup
    card_x = Inches(8.2)
    card_y = Inches(1.5)
    card_w = Inches(4.8)
    card_h = Inches(5.5)
    _add_rect(slide, card_x, card_y, card_w, card_h, fill=SURFACE, line=BORDER)

    # Header inside card
    _add_text(slide, card_x + Inches(0.2), card_y + Inches(0.15), card_w - Inches(0.4),
              Inches(0.3), "INC-1042", size=13, bold=True, color=TEXT)
    pill = _add_rect(slide, card_x + card_w - Inches(1.3), card_y + Inches(0.15),
                     Inches(1.1), Inches(0.3), fill=AMBER_BG)
    _add_text(slide, card_x + card_w - Inches(1.3), card_y + Inches(0.15),
              Inches(1.1), Inches(0.3), "In progress", size=10, bold=True,
              color=AMBER, align=PP_ALIGN.CENTER)

    # Run-check block (dashed accent border feel via teal line)
    rb_x = card_x + Inches(0.2)
    rb_y = card_y + Inches(0.65)
    rb_w = card_w - Inches(0.4)
    rb_h = Inches(1.6)
    block = _add_rect(slide, rb_x, rb_y, rb_w, rb_h, fill=ACCENT_BG, line=TEAL)
    _add_text(slide, rb_x + Inches(0.15), rb_y + Inches(0.1), rb_w - Inches(0.3),
              Inches(0.3), "Run check: Lead status", size=11, bold=True, color=TEAL)
    _add_text(slide, rb_x + Inches(0.15), rb_y + Inches(0.45), Inches(0.8),
              Inches(0.3), "Lead ID", size=10, color=MUTED)
    input_box = _add_rect(slide, rb_x + Inches(0.15), rb_y + Inches(0.7),
                          rb_w - Inches(0.3), Inches(0.35), fill=SURFACE, line=BORDER)
    _add_text(slide, rb_x + Inches(0.2), rb_y + Inches(0.73), rb_w - Inches(0.4),
              Inches(0.3), "Enter your Lead ID", size=10, color=MUTED)
    _add_rect(slide, rb_x + Inches(0.15), rb_y + Inches(1.1), rb_w - Inches(0.3),
              Inches(0.4), fill=NAVY)
    _add_text(slide, rb_x + Inches(0.15), rb_y + Inches(1.15), rb_w - Inches(0.3),
              Inches(0.3), "Run", size=11, bold=True, color=SURFACE, align=PP_ALIGN.CENTER)

    # Result table
    rt_y = rb_y + rb_h + Inches(0.2)
    _add_text(slide, rb_x, rt_y, rb_w, Inches(0.3), "Result", size=11, bold=True, color=TEXT)
    _add_rect(slide, rb_x, rt_y + Inches(0.35), rb_w, Inches(0.05), fill=BORDER)  # header sep
    _add_text(slide, rb_x + Inches(0.1), rt_y + Inches(0.4), Inches(1.5),
              Inches(0.3), "Field", size=10, bold=True, color=MUTED)
    _add_text(slide, rb_x + rb_w / 2, rt_y + Inches(0.4), rb_w / 2,
              Inches(0.3), "Value", size=10, bold=True, color=MUTED)
    rows = [("status", "Approved"), ("disbursementDate", "2026-08-12")]
    for i, (k, v) in enumerate(rows):
        y = rt_y + Inches(0.7) + Inches(0.35) * i
        _add_text(slide, rb_x + Inches(0.1), y, Inches(1.5), Inches(0.3), k, size=10, color=TEXT)
        _add_text(slide, rb_x + rb_w / 2, y, rb_w / 2 - Inches(0.1), Inches(0.3), v, size=10, color=TEXT)

    # Counter line
    _add_text(slide, rb_x, rt_y + Inches(1.5), rb_w, Inches(0.3),
              "You've run 4 checks on this ticket", size=9, color=MUTED)

    # Footer note
    _add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.3),
              "Per-ticket counter increments at run start — failed validations still count, "
              "so the reporter always sees the tally move the moment they press Run.",
              size=11, color=MUTED)


def main():
    if not os.path.exists(DECK):
        print(f"Design deck not found: {DECK}", file=sys.stderr)
        sys.exit(1)
    prs = Presentation(DECK)
    add_slide(prs)
    prs.save(DECK)
    print(f"Added slide; deck now has {len(prs.slides)} slides.")


if __name__ == "__main__":
    main()
