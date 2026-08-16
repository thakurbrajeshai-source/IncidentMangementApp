"""
Reporting job for the Incident Management App.

Reads the same DB the .NET API uses (SQLite for dev, MS SQL for prod), then
generates the two files that the Admin "Generate report" button delivers:

  1. Admin_Report_YYYYMMDD_HHMMSS.xlsx   (6 sheets, matches the validated template)
  2. Admin_Report_YYYYMMDD_HHMMSS.pptx   (executive deck, one chart per slide)

Sheet structure (matches the template in
D:\\IncidentMangementApp\\Admin_Report_Format_Template.xlsx):
  - Incident Log
  - Resolution Evolution Matrix  (category x month)
  - Reporter Leaderboard
  - Resolver Leaderboard
  - Status Summary
  - Rejection Log

Triggered by ReportsController.Generate() in the .NET backend, which invokes
this script as a subprocess and serves the resulting files back to the admin.

Usage:
  python generate_reports.py
  DB_PROVIDER=Sqlite  DB_CONNECTION_STRING="Data Source=.../incident_management.db"  python generate_reports.py
  DB_PROVIDER=SqlServer DB_CONNECTION_STRING="Driver={ODBC Driver 17 for SQL Server};Server=...;Database=IncidentManagement;Trusted_Connection=yes;"  python generate_reports.py
"""

import os
import sys
import sqlite3
from datetime import datetime
from pathlib import Path

import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, Reference
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# ---- DB connection ---------------------------------------------------------

DB_PROVIDER = os.environ.get("DB_PROVIDER", "Sqlite")
DB_CONNECTION_STRING = os.environ.get("DB_CONNECTION_STRING",
    # Default: look for the dev SQLite file the .NET API created
    str(Path(__file__).parent.parent / "backend" / "incident_management.db"))


def query(sql: str) -> pd.DataFrame:
    """Run a query and return a DataFrame. Works for both SQLite and MS SQL."""
    if DB_PROVIDER.lower() == "sqlite":
        conn = sqlite3.connect(DB_CONNECTION_STRING)
        try:
            return pd.read_sql_query(sql, conn)
        finally:
            conn.close()
    elif DB_PROVIDER.lower() in ("sqlserver", "mssql"):
        # Production path. Uses pyodbc directly (avoids a SQLAlchemy dep for a
        # one-shot batch job). The .NET backend passes the connection string
        # with the right ODBC driver.
        import pyodbc
        conn = pyodbc.connect(DB_CONNECTION_STRING)
        try:
            return pd.read_sql_query(sql, conn)
        finally:
            conn.close()
    else:
        raise SystemExit(f"Unknown DB_PROVIDER: {DB_PROVIDER}")


# ---- SQL: the same shape regardless of provider, but we map enum ints to
# the same names the .NET domain uses. EF stores enums as ints.

# (matching Domain/Enums.cs)
ENUM_USER_ROLE = {1: "Reporter", 2: "Resolver", 3: "Admin"}
ENUM_USER_STATUS = {1: "Active", 2: "Disabled"}
ENUM_INCIDENT_STATUS = {1: "Open", 2: "InProgress", 3: "Resolved", 4: "Closed", 5: "Rejected", 6: "Reopened"}
ENUM_ASSIGNMENT_TYPE = {1: "SelfPicked", 2: "AdminAssigned", 3: "Tagged", 4: "Reassigned"}


def load_data() -> dict[str, pd.DataFrame]:
    """Pull all data we need into a dict of DataFrames."""
    users = query("""
        SELECT Id, Mobile, FirstName, LastName, Email, Role, Status, CreatedAt
        FROM Users
    """)
    categories = query("SELECT Id, Name, Description FROM Categories ORDER BY Id")
    incidents = query("""
        SELECT i.Id, i.TicketRef, i.ReporterId, i.CategoryId, c.Name AS CategoryName,
               i.Description, i.Status, i.CurrentAssigneeId, i.RejectionReason,
               i.RejectedById, i.CreatedAt, i.ResolvedAt, i.ClosedAt, i.RejectedAt, i.RevertCount
        FROM Incidents i
        JOIN Categories c ON c.Id = i.CategoryId
    """)
    assignments = query("""
        SELECT a.Id, a.IncidentId, a.ResolverId, a.AssignmentType, a.AssignedAt
        FROM IncidentAssignments a
    """)
    return {
        "users": users,
        "categories": categories,
        "incidents": incidents,
        "assignments": assignments,
    }


# ---- Formatting helpers ----------------------------------------------------

NAVY = "1E3A8A"
TEAL = "0D9488"
MINT = "10B981"
AMBER = "D97706"
DANGER = "B91C1C"
MUTED = "6B7280"
BG = "F9FAFB"
BORDER = "E5E7EB"

HEADER_FILL = PatternFill("solid", fgColor=NAVY)
HEADER_FONT = Font(color="FFFFFF", bold=True, size=11)
THIN = Side(border_style="thin", color=BORDER)
BORDER_STYLE = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)


def style_header(ws, row: int, cols: int):
    for c in range(1, cols + 1):
        cell = ws.cell(row=row, column=c)
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
        cell.alignment = Alignment(horizontal="left", vertical="center")
    ws.row_dimensions[row].height = 20


def autosize(ws, max_width: int = 40):
    for col in ws.columns:
        col_letter = get_column_letter(col[0].column)
        max_len = max((len(str(c.value)) for c in col if c.value is not None), default=10)
        ws.column_dimensions[col_letter].width = min(max_len + 2, max_width)


def fmt_mobile(m):
    """Format a 10-digit mobile as +91 XXXXX XXXXX for display."""
    if m is None or pd.isna(m): return ""
    s = str(m).strip()
    if len(s) == 10 and s.isdigit(): return f"+91 {s[:5]} {s[5:]}"
    return s


def name_of(user_id, users: pd.DataFrame) -> str:
    if user_id is None or pd.isna(user_id): return ""
    u = users[users["Id"] == user_id]
    if u.empty: return ""
    r = u.iloc[0]
    return f"{r['FirstName']} {r['LastName']}".strip()


# ---- Excel generation ------------------------------------------------------

def build_excel(data: dict, out_path: Path):
    users = data["users"]
    incidents = data["incidents"].copy()
    assignments = data["assignments"].copy()

    # Decode enums (stored as int in EF)
    incidents["StatusLabel"] = incidents["Status"].map(ENUM_INCIDENT_STATUS)
    users["RoleLabel"] = users["Role"].map(ENUM_USER_ROLE)
    assignments["AssignmentLabel"] = assignments["AssignmentType"].map(ENUM_ASSIGNMENT_TYPE)

    # Helpful lookups
    incidents["ReporterMobile"] = incidents["ReporterId"].apply(lambda i: fmt_mobile(users.loc[users["Id"] == i, "Mobile"].squeeze() if (users["Id"] == i).any() else None))
    incidents["AssigneeName"]    = incidents["CurrentAssigneeId"].apply(lambda i: name_of(i, users))
    incidents["RejectedByName"]  = incidents["RejectedById"].apply(lambda i: name_of(i, users))
    incidents["Category"]        = incidents["CategoryName"]

    wb = Workbook()
    wb.remove(wb.active)

    # ----- Sheet: Incident Log ---------------------------------------------
    ws = wb.create_sheet("Incident Log")
    headers = ["Ticket Ref", "Category", "Reporter", "Assignee (Resolver)", "Status",
               "Created At", "Resolved At", "Closed At", "Revert Count", "Rejection Reason", "Rejected By"]
    ws.append(headers); style_header(ws, 1, len(headers))
    for _, r in incidents.sort_values("CreatedAt", ascending=False).iterrows():
        ws.append([
            r["TicketRef"], r["Category"], r["ReporterMobile"],
            r["AssigneeName"] or "Unassigned",
            r["StatusLabel"] if r["StatusLabel"] != "Reopened" else "In Progress (reverted)",
            pd.to_datetime(r["CreatedAt"]).strftime("%Y-%m-%d %H:%M") if pd.notna(r["CreatedAt"]) else "",
            pd.to_datetime(r["ResolvedAt"]).strftime("%Y-%m-%d %H:%M") if pd.notna(r["ResolvedAt"]) else "",
            pd.to_datetime(r["ClosedAt"]).strftime("%Y-%m-%d %H:%M") if pd.notna(r["ClosedAt"]) else "",
            int(r["RevertCount"]) if pd.notna(r["RevertCount"]) else 0,
            r["RejectionReason"] or "",
            r["RejectedByName"] or "",
        ])
    autosize(ws)
    ws.freeze_panes = "A2"

    # ----- Sheet: Resolution Evolution Matrix ------------------------------
    ws = wb.create_sheet("Resolution Evolution Matrix")
    ws.cell(row=1, column=1, value="Category x Month Evolution Matrix").font = Font(bold=True, size=13)
    ws.cell(row=2, column=1, value="Incident volume per category over time, from live ticket data.").font = Font(italic=True, color=MUTED, size=10)

    df = incidents.copy()
    df["CreatedMonth"] = pd.to_datetime(df["CreatedAt"]).dt.to_period("M").astype(str)
    pivot = df.pivot_table(index="Category", columns="CreatedMonth", values="Id", aggfunc="count", fill_value=0)
    pivot = pivot.sort_index()
    months = list(pivot.columns)

    # Header
    ws.cell(row=4, column=1, value="Category")
    for i, m in enumerate(months, start=2):
        ws.cell(row=4, column=i, value=m)
    style_header(ws, 4, 1 + len(months))

    for ri, (cat, row) in enumerate(pivot.iterrows(), start=5):
        ws.cell(row=ri, column=1, value=cat)
        for ci, m in enumerate(months, start=2):
            ws.cell(row=ri, column=ci, value=int(row[m]))
    autosize(ws, max_width=25)

    # Bar chart of the latest month by category
    chart = BarChart()
    chart.type = "bar"
    chart.style = 11
    chart.title = "Incidents by category (latest month)"
    chart.y_axis.title = "Category"
    chart.x_axis.title = "Count"
    if months:
        last_col = 1 + len(months)
        data_ref = Reference(ws, min_col=last_col, min_row=4, max_col=last_col, max_row=4 + len(pivot))
        cat_ref  = Reference(ws, min_col=1, min_row=5, max_row=4 + len(pivot))
        chart.add_data(data_ref, titles_from_data=True)
        chart.set_categories(cat_ref)
        chart.width = 18; chart.height = 9
        ws.add_chart(chart, f"A{4 + len(pivot) + 3}")

    # ----- Sheet: Reporter Leaderboard -------------------------------------
    ws = wb.create_sheet("Reporter Leaderboard")
    headers = ["Reporter", "Incidents Reported", "Currently Open", "Currently Closed", "Reverted (fix rejected)"]
    ws.append(headers); style_header(ws, 1, len(headers))
    reportable = incidents[incidents["StatusLabel"].isin(["Open", "InProgress", "Resolved", "Closed", "Reopened"])]
    rep = reportable.groupby("ReporterMobile").agg(
        Reported=("Id", "count"),
        Open=("StatusLabel", lambda s: (s.isin(["Open", "InProgress", "Reopened"])).sum()),
        Closed=("StatusLabel", lambda s: (s == "Closed").sum()),
        Reverted=("RevertCount", "sum"),
    ).reset_index().sort_values(["Reported", "Reverted"], ascending=[False, False])
    for _, r in rep.iterrows():
        ws.append([r["ReporterMobile"], int(r["Reported"]), int(r["Open"]), int(r["Closed"]), int(r["Reverted"])])
    autosize(ws, max_width=30)

    # ----- Sheet: Resolver Leaderboard -------------------------------------
    ws = wb.create_sheet("Resolver Leaderboard")
    headers = ["Resolver", "Tickets Picked/Assigned", "Resolved", "Avg. Resolution Time (hrs)", "Currently Active"]
    ws.append(headers); style_header(ws, 1, len(headers))

    # Picked/Assigned = distinct incidents in IncidentAssignments (any type)
    picked = assignments.groupby("ResolverId")["IncidentId"].nunique().rename("Picked")
    resolved = incidents[incidents["StatusLabel"].isin(["Resolved", "Closed", "Reopened"])].groupby("CurrentAssigneeId")["Id"].count().rename("Resolved")
    active = incidents[incidents["StatusLabel"].isin(["InProgress", "Reopened", "Resolved"])].groupby("CurrentAssigneeId")["Id"].count().rename("Active")

    # Avg resolution time (hrs) = avg(ResolvedAt - CreatedAt) where status is Resolved/Closed
    inc_done = incidents[incidents["StatusLabel"].isin(["Resolved", "Closed"])].copy()
    inc_done["Hours"] = (pd.to_datetime(inc_done["ResolvedAt"]) - pd.to_datetime(inc_done["CreatedAt"])).dt.total_seconds() / 3600
    avg_hrs = inc_done.groupby("CurrentAssigneeId")["Hours"].mean().rename("AvgHrs")

    resolvers = users[users["RoleLabel"] == "Resolver"]
    rows = []
    for _, u in resolvers.iterrows():
        uid = u["Id"]
        rows.append([
            f"{u['FirstName']} {u['LastName']}".strip(),
            int(picked.get(uid, 0)),
            int(resolved.get(uid, 0)),
            round(float(avg_hrs.get(uid, 0)) if pd.notna(avg_hrs.get(uid, 0)) else 0, 1),
            int(active.get(uid, 0)),
        ])
    rows.sort(key=lambda r: r[1], reverse=True)
    for row in rows: ws.append(row)
    autosize(ws, max_width=30)

    # ----- Sheet: Status Summary -------------------------------------------
    ws = wb.create_sheet("Status Summary")
    headers = ["Status", "Ticket Count", "Avg. Time in State (hrs)"]
    ws.append(headers); style_header(ws, 1, len(headers))

    # Count by status, with Reopened collapsed into InProgress
    inc_for_summary = incidents.copy()
    inc_for_summary["DisplayStatus"] = inc_for_summary["StatusLabel"].replace({"Reopened": "InProgress"})
    status_counts = inc_for_summary["DisplayStatus"].value_counts().to_dict()

    # For the Open / InProgress rows, avg time = avg(now - CreatedAt) (current age).
    # For Resolved, avg time = avg(ResolvedAt - CreatedAt). Closed/Rejected are terminal.
    def avg_hours_since(series_dates, baseline_dates):
        d = (pd.to_datetime(series_dates) - pd.to_datetime(baseline_dates)).dt.total_seconds() / 3600
        return d[d >= 0].mean() if len(d) else 0

    now = pd.Timestamp.utcnow().tz_localize(None)  # match tz-naive strings from SQLite/SQL
    open_avg    = avg_hours_since(inc_for_summary.loc[inc_for_summary["DisplayStatus"] == "Open", "CreatedAt"],
                                  inc_for_summary.loc[inc_for_summary["DisplayStatus"] == "Open", "CreatedAt"])  # age = now - created; we approximate with sample
    # Simpler: just avg of (now - createdAt) for currently-Open / InProgress rows
    def avg_age(filter_status):
        sub = inc_for_summary[inc_for_summary["DisplayStatus"] == filter_status]
        if sub.empty: return 0
        return float((now - pd.to_datetime(sub["CreatedAt"])).dt.total_seconds().mean() / 3600)

    rows = [
        ["Open",                          int(status_counts.get("Open", 0)),       round(avg_age("Open"), 1)],
        ["In Progress",                   int(status_counts.get("InProgress", 0) + status_counts.get("Reopened", 0)), round(avg_age("InProgress"), 1)],
        ["Resolved (awaiting confirm)",   int(status_counts.get("Resolved", 0)),  round(avg_hours_since(inc_for_summary.loc[inc_for_summary["StatusLabel"] == "Resolved", "ResolvedAt"],
                                                                                       inc_for_summary.loc[inc_for_summary["StatusLabel"] == "Resolved", "CreatedAt"]), 1)],
        ["Closed",                        int(status_counts.get("Closed", 0)),    ""],
        ["Rejected",                      int(status_counts.get("Rejected", 0)),  ""],
        ["Reverted at least once",        int((incidents["RevertCount"] > 0).sum()), ""],
    ]
    for row in rows: ws.append(row)
    autosize(ws)

    # Bar chart of status counts
    chart = BarChart()
    chart.type = "col"
    chart.style = 11
    chart.title = "Ticket count by status"
    chart.y_axis.title = "Count"
    chart.x_axis.title = "Status"
    data_ref = Reference(ws, min_col=2, min_row=1, max_col=2, max_row=1 + len(rows))
    cat_ref  = Reference(ws, min_col=1, min_row=2, max_row=1 + len(rows))
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cat_ref)
    chart.width = 16; chart.height = 9
    ws.add_chart(chart, "E2")

    # ----- Sheet: Rejection Log --------------------------------------------
    ws = wb.create_sheet("Rejection Log")
    headers = ["Ticket Ref", "Reporter", "Category", "Rejected By (Admin)", "Reason", "Rejected At"]
    ws.append(headers); style_header(ws, 1, len(headers))
    rejected = incidents[incidents["StatusLabel"] == "Rejected"].copy()
    for _, r in rejected.sort_values("CreatedAt", ascending=False).iterrows():
        ws.append([
            r["TicketRef"], r["ReporterMobile"], r["Category"],
            r["RejectedByName"] or "Admin",
            r["RejectionReason"] or "", pd.to_datetime(r["RejectedAt"]).strftime("%Y-%m-%d %H:%M") if pd.notna(r["RejectedAt"]) else "",
        ])
    autosize(ws)

    wb.save(out_path)


# ---- PowerPoint generation -------------------------------------------------

def build_pptx(data: dict, out_path: Path, status_counts: dict):
    incidents = data["incidents"].copy()
    incidents["StatusLabel"] = incidents["Status"].map(ENUM_INCIDENT_STATUS)
    categories = data["categories"]
    users = data["users"]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_slide(title, subtitle=None):
        s = prs.slides.add_slide(blank)
        # Header bar
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
        bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(NAVY)
        bar.line.fill.background()
        tx = s.shapes.add_textbox(Inches(0.5), Inches(0.15), prs.slide_width - Inches(1), Inches(0.6))
        tf = tx.text_frame; tf.text = title
        p = tf.paragraphs[0]; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = RGBColor.from_string("FFFFFF")
        if subtitle:
            sub = s.shapes.add_textbox(Inches(0.5), Inches(1.0), prs.slide_width - Inches(1), Inches(0.4))
            sf = sub.text_frame; sf.text = subtitle
            sp = sf.paragraphs[0]; sp.font.size = Pt(12); sp.font.color.rgb = RGBColor.from_string(MUTED)
        return s

    def add_bullets(slide, items, x=Inches(0.5), y=Inches(1.6), w=Inches(12.3), h=Inches(5.0)):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        for i, t in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + t
            p.font.size = Pt(14); p.font.color.rgb = RGBColor.from_string("111827")
            p.space_after = Pt(6)

    def add_kpi_row(slide, kpis, y=Inches(2.0)):
        for i, (label, value, color) in enumerate(kpis):
            x = Inches(0.5 + i * 3.2)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.0), Inches(1.2))
            box.fill.solid(); box.fill.fore_color.rgb = RGBColor.from_string(BG)
            box.line.color.rgb = RGBColor.from_string(BORDER)
            tb = slide.shapes.add_textbox(x, y + Inches(0.15), Inches(3.0), Inches(0.9))
            tf = tb.text_frame
            tf.text = str(value)
            tf.paragraphs[0].font.size = Pt(28); tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor.from_string(color)
            p2 = tf.add_paragraph(); p2.text = label
            p2.font.size = Pt(11); p2.font.color.rgb = RGBColor.from_string(MUTED)

    def add_chart(slide, df, x, y, w, h, kind="bar", title=None):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        cdata = CategoryChartData()
        cdata.categories = list(df.iloc[:, 0])
        cdata.add_series("Count", list(df.iloc[:, 1]))
        ct = XL_CHART_TYPE.BAR_CLUSTERED if kind == "bar" else XL_CHART_TYPE.COLUMN_CLUSTERED
        chart = slide.shapes.add_chart(ct, x, y, w, h, cdata).chart
        chart.has_title = bool(title)
        if title: chart.chart_title.text_frame.text = title
        return chart

    total = len(incidents)

    # 1. Title
    add_slide("Incident Management — Executive Summary",
              f"Auto-generated report · {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}")
    add_bullets(prs.slides[-1], [
        f"Total tickets in system: {total}",
        f"Open: {status_counts.get('Open', 0)} · In Progress: {status_counts.get('InProgress', 0) + status_counts.get('Reopened', 0)} · "
        f"Resolved: {status_counts.get('Resolved', 0)} · Closed: {status_counts.get('Closed', 0)} · Rejected: {status_counts.get('Rejected', 0)}",
        f"Reverted at least once: {(incidents['RevertCount'] > 0).sum()}",
        "This deck is generated live from the Incident Management App database.",
    ])

    # 2. KPI cards
    s = add_slide("Key metrics")
    add_kpi_row(s, [
        ("Open",        status_counts.get("Open", 0),                                          AMBER),
        ("In Progress", status_counts.get("InProgress", 0) + status_counts.get("Reopened", 0), TEAL),
        ("Closed",      status_counts.get("Closed", 0),                                        MINT),
        ("Reverted",    int((incidents["RevertCount"] > 0).sum()),                             DANGER),
    ], y=Inches(2.5))

    # 3. Monthly trend (by CreatedAt)
    s = add_slide("Monthly incident trend", "Ticket creation volume by month")
    inc = incidents.copy()
    inc["Month"] = pd.to_datetime(inc["CreatedAt"]).dt.to_period("M").astype(str)
    monthly = inc.groupby("Month").size().reset_index(name="Count")
    add_chart(s, monthly, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.0), kind="col")

    # 4. Category breakdown
    s = add_slide("Category breakdown", "Where do incidents cluster?")
    cat = inc.groupby("CategoryName").size().reset_index(name="Count").sort_values("Count", ascending=False)
    add_chart(s, cat, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.0), kind="bar")

    # 5. Resolver leaderboard
    s = add_slide("Resolver leaderboard", "Tickets picked/assigned per resolver")
    assignments = data["assignments"].copy()
    picked = assignments.groupby("ResolverId")["IncidentId"].nunique().reset_index(name="Picked")
    picked["Resolver"] = picked["ResolverId"].apply(lambda i: name_of(i, users))
    picked = picked[["Resolver", "Picked"]].sort_values("Picked", ascending=False).head(10)
    add_chart(s, picked, Inches(0.5), Inches(1.8), Inches(12.3), Inches(5.0), kind="bar")

    # 6. Key findings (auto-derived)
    s = add_slide("Key findings")
    top_cat = cat.iloc[0]["CategoryName"] if not cat.empty else "n/a"
    top_res = picked.iloc[0]["Resolver"] if not picked.empty else "n/a"
    revert_pct = (incidents["RevertCount"] > 0).mean() * 100
    add_bullets(s, [
        f"Highest-volume category: {top_cat}",
        f"Most active resolver: {top_res}",
        f"Revert rate (tickets reopened at least once): {revert_pct:.1f}%",
        f"Total rejections recorded: {status_counts.get('Rejected', 0)}",
        "Recommendations: focus on the top category, watch reverts for training signals, and review rejections for policy gaps.",
    ])

    prs.save(out_path)


# ---- Entry point -----------------------------------------------------------

def main():
    out_dir = Path(__file__).parent / "output"
    out_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    xlsx_path = out_dir / f"Admin_Report_{stamp}.xlsx"
    pptx_path = out_dir / f"Admin_Report_{stamp}.pptx"

    print(f"Reading from {DB_PROVIDER} db: {DB_CONNECTION_STRING}")
    data = load_data()
    print(f"Loaded {len(data['users'])} users, {len(data['incidents'])} incidents, {len(data['assignments'])} assignments")

    build_excel(data, xlsx_path)
    print(f"Wrote {xlsx_path}")

    status_counts = data["incidents"]["Status"].map(ENUM_INCIDENT_STATUS).value_counts().to_dict()
    build_pptx(data, pptx_path, status_counts)
    print(f"Wrote {pptx_path}")


if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"FAILED: {e}", file=sys.stderr)
        sys.exit(1)
