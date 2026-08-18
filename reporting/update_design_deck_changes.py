"""
Adds a 'Changes Applied — PRD Sync' slide to the Design Deck documenting
all features implemented during the current development cycle.

Run from repo root:
    python src/reporting/update_design_deck_changes.py
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

DECK = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "..", "..", "Incident_App_Design_Deck.pptx",
)
DECK = os.path.normpath(DECK)

NAVY = RGBColor(0x1E, 0x3A, 0x8A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
MUTED = RGBColor(0x6B, 0x72, 0x80)
TEXT = RGBColor(0x11, 0x18, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
GREEN = RGBColor(0x05, 0x96, 0x69)
AMBER_BG = RGBColor(0xFE, 0xF3, 0xC7)
LIGHT_BG = RGBColor(0xF9, 0xFA, 0xFB)


def _add_text(slide, left, top, width, height, text, *, size=14, bold=False,
              color=TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_rect(slide, left, top, width, height, *, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def _add_badge(slide, left, top, text, *, fill=GREEN_BG, text_color=GREEN):
    """Small rounded badge pill."""
    w, h = Inches(1.6), Inches(0.28)
    _add_rect(slide, left, top, w, h, fill=fill, line=None)
    _add_text(slide, left, top, w, h, text, size=8, bold=True,
              color=text_color, align=PP_ALIGN.CENTER)
    return w


def add_slide(prs):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # ── Title ──────────────────────────────────────────────────────────────
    _add_text(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
              "Changes Applied — PRD Sync", size=26, bold=True, color=NAVY)
    _add_text(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.35),
              "All features from the current development cycle, verified and integrated",
              size=13, color=MUTED)

    # ── Column 1: Backend Features ─────────────────────────────────────────
    col1_x = Inches(0.5)
    col1_w = Inches(5.8)
    y = Inches(1.3)

    _add_text(slide, col1_x, y, col1_w, Inches(0.35),
              "Backend Features", size=16, bold=True, color=NAVY)
    y += Inches(0.4)

    features = [
        ("WorkflowCategory join table", "Links workflows to categories as 'default check'; composite PK"),
        ("Default-workflow API", "GET /api/incidents/default-workflow/{categoryId}"),
        ("Attach workflow to ticket", "POST /api/workflows/attach + WorkflowIncidentAssignment entity"),
        ("Per-user run counter", "WorkflowRunCounter incremented at run start, exposed via /run-count"),
        ("Reporter trigger", "POST /api/incidents/{id}/run-workflow (reporter-only on own tickets)"),
        ("Visibility toggle", "PUT /api/workflows/visibility — show/hide output in thread"),
        ("OTP TTL enforcement", "ConcurrentDictionary tracking; skipped in dev/test mode"),
        ("CORS configuration", "Cors:AllowedOrigins in appsettings; HTTPS redirect in prod"),
    ]
    for name, desc in features:
        _add_rect(slide, col1_x, y, col1_w, Inches(0.55), fill=LIGHT_BG, line=BORDER)
        _add_badge(slide, col1_x + Inches(0.1), y + Inches(0.13), "DONE")
        _add_text(slide, col1_x + Inches(1.8), y + Inches(0.05), col1_w - Inches(2.0),
                  Inches(0.22), name, size=11, bold=True, color=TEXT)
        _add_text(slide, col1_x + Inches(1.8), y + Inches(0.28), col1_w - Inches(2.0),
                  Inches(0.22), desc, size=9, color=MUTED)
        y += Inches(0.62)

    # ── Column 2: Frontend Features ────────────────────────────────────────
    col2_x = Inches(6.7)
    col2_w = Inches(5.8)
    y = Inches(1.3)

    _add_text(slide, col2_x, y, col2_w, Inches(0.35),
              "Frontend Features", size=16, bold=True, color=NAVY)
    y += Inches(0.4)

    features2 = [
        ("ReporterWorkflowCard", "Default-workflow card on reporter ticket page"),
        ("AttachWorkflowPicker", "Dropdown to pick + run a workflow; attached to thread"),
        ("Workflow category chips", "Multi-select checkboxes in BuilderModal + table column"),
        ("Visibility toggle UI", "Checkbox per run in thread; Resolver/Admin only"),
        ("Run counter display", "'You've run N checks on this ticket' header"),
        ("MentionInput component", "Shared; outside-click dismiss; inline removable tag chips"),
        ("api.ts types + methods", "AttachedWorkflow, DefaultWorkflow, AvailableWorkflow, RunCountResponse"),
        ("Workflows.tsx updates", "Category multi-select in BuilderModal + Categories table column"),
    ]
    for name, desc in features2:
        _add_rect(slide, col2_x, y, col2_w, Inches(0.55), fill=LIGHT_BG, line=BORDER)
        _add_badge(slide, col2_x + Inches(0.1), y + Inches(0.13), "DONE")
        _add_text(slide, col2_x + Inches(1.8), y + Inches(0.05), col2_w - Inches(2.0),
                  Inches(0.22), name, size=11, bold=True, color=TEXT)
        _add_text(slide, col2_x + Inches(1.8), y + Inches(0.28), col2_w - Inches(2.0),
                  Inches(0.22), desc, size=9, color=MUTED)
        y += Inches(0.62)

    # ── Bottom: Migrations & Infrastructure ────────────────────────────────
    bot_y = Inches(6.5)
    _add_rect(slide, Inches(0.5), bot_y, Inches(12), Inches(0.9), fill=AMBER_BG, line=None)
    _add_text(slide, Inches(0.7), bot_y + Inches(0.08), Inches(11.5), Inches(0.28),
              "Migrations & Infrastructure", size=13, bold=True, color=NAVY)
    _add_text(slide, Inches(0.7), bot_y + Inches(0.35), Inches(11.5), Inches(0.45),
              "EF Migrations: AddWorkflowCategoriesAndRunCounters, AddWorkflowIncidentAssignments  |  "
              "IncidentManagement.sln created  |  "
              "pyodbc + SQLAlchemy in requirements.txt  |  "
              "Key rotation documented in appsettings.json",
              size=10, color=TEXT)


def main():
    if not os.path.exists(DECK):
        print(f"Design deck not found: {DECK}", file=sys.stderr)
        sys.exit(1)
    prs = Presentation(DECK)
    add_slide(prs)
    prs.save(DECK)
    print(f"Added 'Changes Applied' slide; deck now has {len(prs.slides)} slides.")


if __name__ == "__main__":
    main()
